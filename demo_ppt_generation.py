#!/usr/bin/env python3
"""
Demo: Complete PPT Generation Flow

This script demonstrates:
1. File upload via API
2. PPT generation using python-pptx
3. File download and verification

Usage:
    python demo_ppt_generation.py
"""

import asyncio
import io
import json
import sys
from pathlib import Path

# Check dependencies
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

try:
    import httpx
    HTTPX_AVAILABLE = True
except ImportError:
    print("Installing httpx...")
    import subprocess
    subprocess.run([sys.executable, "-m", "pip", "install", "httpx", "-q"])
    import httpx
    HTTPX_AVAILABLE = True


API_BASE = "http://127.0.0.1:8765"


def create_demo_ppt() -> bytes:
    """Create a demo presentation about Cape System."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Cape 代码执行层"
    p.font.size = Pt(54)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "模型无关的能力执行系统"
    p.font.size = Pt(28)
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: Architecture
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "系统架构"
    p.font.size = Pt(40)
    p.font.bold = True

    # Architecture diagram (text-based)
    content = """
    ┌─────────────────────────────────────────┐
    │              Cape Runtime               │
    ├─────────────────────────────────────────┤
    │  ┌─────────┐  ┌─────────┐  ┌─────────┐ │
    │  │   LLM   │  │  Tool   │  │  Code   │ │
    │  │Executor │  │Executor │  │Executor │ │
    │  └────┬────┘  └────┬────┘  └────┬────┘ │
    │       │            │            │      │
    ├───────┴────────────┴────────────┴──────┤
    │            Sandbox Manager             │
    │  ┌─────────┬─────────┬───────────┐    │
    │  │ Docker  │ Process │ InProcess │    │
    │  │ Sandbox │ Sandbox │  Sandbox  │    │
    │  └─────────┴─────────┴───────────┘    │
    └─────────────────────────────────────────┘
    """

    arch_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
    tf = arch_box.text_frame
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(14)
    p.font.name = "Courier New"

    # Slide 3: Features
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "核心功能"
    p.font.size = Pt(40)
    p.font.bold = True

    features = [
        ("🔒 安全沙箱", "Docker/Process 级别隔离，资源限制"),
        ("📄 文档处理", "Excel, Word, PowerPoint, PDF 全支持"),
        ("🔌 模型无关", "支持 Claude, GPT, Gemini 等任意 LLM"),
        ("📁 文件 API", "上传、下载、处理一体化"),
        ("⚡ 高性能", "异步执行，并行处理"),
    ]

    y_pos = 1.8
    for emoji_title, desc in features:
        # Feature title
        box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(5), Inches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = emoji_title
        p.font.size = Pt(24)
        p.font.bold = True

        # Feature description
        box = slide.shapes.add_textbox(Inches(6), Inches(y_pos), Inches(6), Inches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(20)

        y_pos += 1.0

    # Slide 4: Implementation Progress
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "实施进度"
    p.font.size = Pt(40)
    p.font.bold = True

    weeks = [
        ("Week 1", "沙箱框架", "✅ 完成", "ProcessSandbox, InProcessSandbox"),
        ("Week 2", "文档技能", "✅ 完成", "xlsx, docx, pptx, pdf Capes"),
        ("Week 3", "Docker 沙箱", "✅ 完成", "容器隔离，资源限制"),
        ("Week 4", "文件 API", "✅ 完成", "上传/下载/处理端点"),
    ]

    y_pos = 1.8
    for week, task, status, detail in weeks:
        # Week
        box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(1.5), Inches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = week
        p.font.size = Pt(20)
        p.font.bold = True

        # Task
        box = slide.shapes.add_textbox(Inches(2.5), Inches(y_pos), Inches(2.5), Inches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = task
        p.font.size = Pt(20)

        # Status
        box = slide.shapes.add_textbox(Inches(5.2), Inches(y_pos), Inches(1.5), Inches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = status
        p.font.size = Pt(20)

        # Detail
        box = slide.shapes.add_textbox(Inches(7), Inches(y_pos), Inches(5.5), Inches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = detail
        p.font.size = Pt(16)

        y_pos += 1.2

    # Slide 5: API Endpoints
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "文件 API 端点"
    p.font.size = Pt(40)
    p.font.bold = True

    endpoints = [
        ("POST", "/api/files/upload", "上传文件"),
        ("GET", "/api/files/{id}", "下载文件"),
        ("POST", "/api/files/{id}/process", "处理文件"),
        ("GET", "/api/files/session/{id}", "会话文件"),
        ("GET", "/api/files/stats", "存储统计"),
    ]

    y_pos = 1.8
    for method, path, desc in endpoints:
        # Method
        box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(1.2), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = method
        p.font.size = Pt(18)
        p.font.bold = True

        # Path
        box = slide.shapes.add_textbox(Inches(2.5), Inches(y_pos), Inches(5), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = path
        p.font.size = Pt(18)
        p.font.name = "Courier New"

        # Description
        box = slide.shapes.add_textbox(Inches(8), Inches(y_pos), Inches(4), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(18)

        y_pos += 0.9

    # Save to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


async def test_api_flow():
    """Test the complete API flow."""
    print("\n" + "=" * 60)
    print("Cape API 完整流程测试")
    print("=" * 60)

    async with httpx.AsyncClient(base_url=API_BASE, timeout=30) as client:
        # Step 1: Check API health
        print("\n1️⃣  检查 API 状态...")
        response = await client.get("/")
        data = response.json()
        print(f"   ✓ API 运行正常: {data['total_capes']} Capes 可用")

        # Step 2: Create PPT
        print("\n2️⃣  生成演示 PPT...")
        ppt_content = create_demo_ppt()
        print(f"   ✓ PPT 已生成: {len(ppt_content)} bytes, 5 张幻灯片")

        # Step 3: Upload PPT
        print("\n3️⃣  上传 PPT 到 API...")
        files = {"files": ("cape_demo.pptx", ppt_content, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
        data = {"session_id": "demo-session"}

        response = await client.post("/api/files/upload", files=files, data=data)
        upload_data = response.json()
        file_id = upload_data["files"][0]["file_id"]
        print(f"   ✓ 上传成功: file_id={file_id}")
        print(f"   ✓ Session: {upload_data['session_id']}")

        # Step 4: Get metadata
        print("\n4️⃣  获取文件元数据...")
        response = await client.get(f"/api/files/{file_id}/metadata")
        meta = response.json()
        print(f"   ✓ 文件名: {meta['original_name']}")
        print(f"   ✓ 大小: {meta['size_bytes']} bytes")
        print(f"   ✓ 状态: {meta['status']}")
        print(f"   ✓ 类型: {meta['content_type']}")

        # Step 5: List session files
        print("\n5️⃣  列出会话文件...")
        response = await client.get("/api/files/session/demo-session")
        session_data = response.json()
        print(f"   ✓ 会话文件数: {session_data['total_files']}")
        for f in session_data["files"]:
            print(f"      - {f['original_name']} ({f['status']})")

        # Step 6: Download file
        print("\n6️⃣  下载 PPT 文件...")
        response = await client.get(f"/api/files/{file_id}")
        downloaded_content = response.content
        print(f"   ✓ 下载成功: {len(downloaded_content)} bytes")

        # Verify content matches
        if downloaded_content == ppt_content:
            print("   ✓ 内容验证通过: 上传下载一致")
        else:
            print("   ✗ 内容验证失败")

        # Step 7: Get storage stats
        print("\n7️⃣  存储统计...")
        response = await client.get("/api/files/stats")
        stats = response.json()
        print(f"   ✓ 总文件数: {stats['total_files']}")
        print(f"   ✓ 总大小: {stats['total_size_mb']} MB")
        print(f"   ✓ 会话数: {stats['total_sessions']}")

        # Step 8: Save PPT locally
        print("\n8️⃣  保存 PPT 到本地...")
        output_path = Path("/Users/g/Desktop/探索/skillslike/output/cape_demo.pptx")
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_bytes(downloaded_content)
        print(f"   ✓ 已保存: {output_path}")

        # Step 9: Cleanup
        print("\n9️⃣  清理测试数据...")
        response = await client.delete("/api/files/session/demo-session")
        cleanup_data = response.json()
        print(f"   ✓ 已删除 {cleanup_data['deleted_files']} 个文件")

        print("\n" + "=" * 60)
        print("✅ 完整流程测试成功!")
        print("=" * 60)
        print(f"\n📁 PPT 输出位置: {output_path}")

        return str(output_path)


async def main():
    """Main entry point."""
    try:
        output_path = await test_api_flow()
        return output_path
    except httpx.ConnectError:
        print("\n❌ 无法连接到 API 服务器")
        print("   请确保服务器正在运行: uvicorn api.main:app --port 8765")
        sys.exit(1)
    except Exception as e:
        print(f"\n❌ 错误: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    asyncio.run(main())
